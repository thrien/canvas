#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
import datetime
import glob
from urllib.request import Request, urlopen
import json
import urllib.parse

# Canvas API
API_URL = "https://umich.instructure.com/api/v1"
TOKEN = ""
COURSE_ID = 850281


def canvas_api(command, parameters=None):
    if not TOKEN:
        raise RuntimeError("No Canvas API access token defined.")
    url = f"{API_URL}/{command}"
    if parameters:
        url += f"?{urllib.parse.urlencode(parameters)}"
    request = Request(url, headers={"Authorization": f"Bearer {TOKEN}"})
    return urlopen(request)


def canvas_quiz_code(lab):
    # the quiz code for each lab is defined in a quiz on Canvas
    with canvas_api(f"courses/{COURSE_ID}/quizzes",
                    parameters={"search_term": f"Quiz {lab:d}:"}) as response:
        quizzes = json.load(response)
    try:
        quiz = next(quiz for quiz in quizzes
                    if quiz["title"].startswith(f"Quiz {lab:d}:"))
    except StopIteration:
        raise RuntimeError(f"Couldn't find quiz for lab {lab:d} on Canvas.")
    return quiz["access_code"]

intros_path = r"C:\\Users\\umthr\\OneDrive - Umich\\Documents\\Teaching\\WN26 PHYSICS 251\\Introductions"
template = intros_path + r"\\Template.pptx"

# Find the latest lab number from existing folders
lab = max(map(lambda s: int(s.removeprefix("lab")), glob.glob("lab??")))

img_path = f"lab{lab:02d}\\groups015.png"

# Load template presentation
prs = Presentation(template)

# Modify title slide
title_slide = prs.slides[0]
subtitle = title_slide.placeholders[1]
subtitle.text = f"Lab {lab:02d} - Section 015"

# Modify group slide
group_slide = prs.slides[1]
# find the existing picture shape
pic_shape = next((shape for shape in group_slide.shapes
                  if shape.shape_type == MSO_SHAPE_TYPE.PICTURE), None)
# replace picture and keep the same position/size
left   = pic_shape.left    # 2487705
top    = pic_shape.top     # 0
width  = pic_shape.width   # 6656295
height = pic_shape.height  # 5143500
group_slide.shapes._spTree.remove(pic_shape._element)
group_slide.shapes.add_picture(img_path, left, top, width=width, height=height)

# Modify quiz slide
if TOKEN:
    quiz_slide = prs.slides[2]
    quiz_code = canvas_quiz_code(lab)
    quiz_slide.placeholders[0].text = quiz_code

# Save the modified presentation
prs.save(f"{intros_path}\\PHYS251 Lab {lab:02d}.pptx")