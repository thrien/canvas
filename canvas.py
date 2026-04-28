#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Utilities for using Canvas as a GSI

This script is intended to be used by GSIs for PHYSICS 151/251 at the
University of Michigan.
"""
# standard library
import os
import argparse
import inspect
from urllib.parse import quote, unquote_plus
from urllib.request import Request, urlopen
#from urllib.error import HTTPError
import json
from random import randrange
# external dependencies
# for sign-in sheets
import numpy as np
import matplotlib.pyplot as plt
# for gradebooks
import pandas as pd
# for introduction slides
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

long_help = False


class VerboseHelpAction(argparse._HelpAction):
    """Help action with optional verbose help message
    
    If this action is invoked with the long --help option it sets a global
    variable `long_help` to True, which can then be used in the help formatter.
    If invoked with the short -h option it sets it to False.
    """
    def __call__(self, parser, namespace, values, option_string=None):
        global long_help
        long_help = True if option_string == "--help" else False

        super().__call__(parser, namespace, values,
                         option_string=option_string)


# configure help message formatting
class CustomHelpFormatter(
        argparse.RawDescriptionHelpFormatter,
        argparse.ArgumentDefaultsHelpFormatter):
    """Help formatter customized to my preferences
    
    - shows available choices for arguments
    - shows default values for arguments
    - preserves formatting of description and epilog
    - limits the width of help messages to 79 characters
    - limits the number of choices in a metavar to 3, ...
      - unless the global variable `long_help` is True
    """

    def __init__(self, *args, width=79, max_choices=3, **kwargs):
        super().__init__(*args, **kwargs)
        self._width = min(width, self._width)
        self._max_choices = max_choices

    def _limit_choices(self, choices):
        if len(choices) > self._max_choices:
            choices[self._max_choices-1:] = ["..."]
        return choices

    def _get_default_metavar_for_choices(self, choices):
        choices = list(choices)
        if not long_help:
            choices = self._limit_choices(choices)
        return "{" + ", ".join(map(str, choices)) + "}"

    def _get_default_metavar_for_positional(self, action):
        if action.choices is not None:
            return self._get_default_metavar_for_choices(action.choices)
        return super()._get_default_metavar_for_positional(action)

    def _get_default_metavar_for_optional(self, action):
        if action.choices is not None:
            return self._get_default_metavar_for_choices(action.choices)
        else:
            return super()._get_default_metavar_for_optional(action)
        
    def _metavar_formatter(self, action, default_metavar):
        if action.metavar is not None:
            result = action.metavar
        else:
            result = default_metavar

        def format(tuple_size):
            if isinstance(result, tuple):
                return result
            else:
                return (result, ) * tuple_size
        return format

    def _get_help_string(self, action):
        help = super()._get_help_string(action)
        if action.choices is not None and '%(choices)' not in help:
            choices = list(action.choices)
            if not long_help:
                choices = self._limit_choices(choices)
            help += f' (choices: {", ".join(map(str, choices))})'
        return help

# Canvas API
API_URL = "https://umich.instructure.com/api/v1"
TOKEN = os.getenv("CANVAS_API_TOKEN")
COURSES = {"PHYS 151 WN25": 734390,
           "PHYS 251 WN26": 850281,
           "PHYS 251 WN26 GSI": 826079}

# Make sure we are in the right directory
os.chdir(os.path.dirname(__file__))

# already existing labs on disk
prefix = "lab"
existing_labs = [int(entry.removeprefix(prefix)) for entry in os.listdir()
                 if os.path.isdir(entry) and entry.startswith(prefix)]
existing_labs.sort()

instructor = ""  # last name, first name
# int for group numbers, 'I' for the instructor (optional), '.' for nothing
table_layout = [[ 1 , 'I', '.'],
                [ 2 , '.',  8 ],
                [ 3 , '.',  7 ],
                [ 4 ,  5 ,  6 ]]
tables = [table for row in table_layout
                for table in row
                if isinstance(table, int)]


def _canvas_api(command, full_url=False, method="GET", parameters={},
                headers={}, data=None):
    """Call the Canvas API with the given command and parameters
    
    For more details on the API see:
    https://developerdocs.instructure.com/services/canvas

    Args:
        command: the API command, e.g., "courses/123/groups"
        full_url: if True, command is treated as a complete URL
        method: HTTP method (GET, PUT, POST, etc.)
        headers: additional HTTP headers
        data: request body data (for PUT/POST)
        verbose: if True, print debug information
    Returns:
        the response as a Python object (list for JSON, else string)
    """
    if not TOKEN:
        raise RuntimeError("No Canvas API access token defined.")

    # see https://developerdocs.instructure.com/services/canvas/oauth2/file.oauth#using-access-tokens
    url = command if full_url else f"{API_URL}/{command}"
    if parameters:
        url += "?" + "&".join(f"{key}={quote(str(value))}"
                              for key, value in parameters.items())
    request = Request(url,
                      headers={"Authorization": f"Bearer {TOKEN}"}|headers,
                      method=method)
    if data is not None:
        request.data = data.encode()

    # Debug output
    if verbose > 2:
        print(f"Request: {method} {url}")
        print(f"Headers: {dict(request.headers)}")
        if data is not None:
            print(f"Data: {data}")

    response = urlopen(request)

    if response.getheader("Content-Type").startswith("application/json"):
        content = json.load(response)
    else:
        content = response.read()

    # see https://developerdocs.instructure.com/services/canvas/basics/file.pagination
    try:  # reading next pages
        links = response.getheader("Link").split(",")
        pages = {rel.removeprefix(" rel=").strip('"'): link.strip("<>")
                 for link, rel in map(lambda page: page.split(";"), links)}
        content += _canvas_api(pages["next"], full_url=True)
    except (KeyError, AttributeError):
        pass

    return content


def _canvas_import_csv(lab):
    # The groups for each lab are defined in a group category on Canvas
    # see https://developerdocs.instructure.com/services/canvas/resources/group_categories#method.group_categories.export
    categories = _canvas_api(f"courses/{COURSE_ID}/group_categories")
    parse_number = lambda name: int(name.lower().removeprefix("lab").strip())
    try:
        category = next(category for category in categories
                        if parse_number(category["name"]) == lab)
    except StopIteration:
        raise RuntimeError(f"Couldn't find groups for lab {lab:d} on Canvas.")

    # Ask Canvas to export the groups for this lab as a CSV
    data = _canvas_api(f"group_categories/{category["id"]}/export")

    # Save CSV on disk
    filename = os.path.join(f"lab{lab:02d}", "canvas.csv")
    os.makedirs(os.path.dirname(filename), exist_ok=True)
    with open(filename, "wb") as file:
        file.write(data)
    if verbose:
        print(f'CSV file written to "{filename}".')


def _format_name(name):
    """Format "Smith, Emma Marie" as "Emma Smith"
    """
    lasts, firsts = name.split(",")
    first = firsts.strip().split(" ")[0]
    last = lasts.strip()  #.split(" ")[0]

    return f"__ {first:s} {last:s}"


def _draw(names, groups, title="Groups", smallfont=18, bigfont=25,
          margin=0.015, title_margin=0.075):
    """Draw the group assignment on tables

    Args:
        names: list of the names of all students
        groups: list of the corresponding group numbers
    Optional arguments:
        title: the title of the figure
        smallfont: font size for names and group numbers
        bigfont: font size for the title
        margin: the margin between tables as a fraction of the figure size
        title_margin: the margin for the title as a fraction of the figure size
    Returns:
        the created matplotlib figure
    """
    # Create a grid of subplots based on table layout
    fig, axes = plt.subplot_mosaic(table_layout, figsize=(11, 8.5))

    fig.suptitle(title, fontsize=bigfont)

    for table in tables:
        # Remove ticks and thicken spine for group axes
        ax = axes[table]
        ax.set_xticks([])
        ax.set_yticks([])
        for spine in ax.spines.values():
            spine.set_linewidth(2)
        # group number in the top left corner
        ax.text(0.05, 0.90, "{:d}".format(table),
                transform=ax.transAxes,
                fontsize=smallfont, ha="left", va="top")
        # the list of names
        ax.text(0.05, 0.70, "\n".join(sorted(names[groups == table])),
                transform=ax.transAxes,
                fontsize=smallfont, ha="left", va="top")

    # Subplot for the instructor (me)
    if 'I' in axes.keys():
        ax = axes['I']
        ax.set_axis_off()
        if instructor:
            ax.text(0.5, 0.5, _format_name(instructor),
                    transform=ax.transAxes,
                    fontsize=smallfont, ha="center", va="center")

    # Adjust spacing between subplots
    fig.subplots_adjust(left=margin, right=1-margin,
                        top=1-title_margin, bottom=margin,
                        wspace=margin*len(table_layout[0]),
                        hspace=margin*len(table_layout))

    return fig


class FlatListAction(argparse.Action):
    """Action for argparse that flattens a list of lists

    With the interval type this parses "-l 1..3 5" into [1, 2, 3, 5].
    """
    def __call__(self, parser, namespace, values, option_string=None):
        setattr(namespace, self.dest, sum(values, start=[]))


def _interval(string, last=max(existing_labs, default=0)):
    """integers or intervals like a..b"""
    wraps = lambda s: last + 1 - int(s[1:]) if s.startswith("-") else int(s)
    if ".." in string:
        a, b = map(wraps, string.split(".."))
        return list(range(a, b + 1))
    else:
        return [wraps(string)]


def _sheets_parser(parser):
    parser.add_argument("-f", "--force", action=argparse.BooleanOptionalAction,
                        help="pull recent CSV from Canvas")
    parser.add_argument("-e", "--extensions", nargs="+",
                        #default=<parsed from function signature>,
                        metavar="ext", help="output formats")
    parser.add_argument("-l", "--labs", type=_interval, nargs="+",
                        action=FlatListAction,
                        default=[max(existing_labs, default=0) + bool(TOKEN)]
                                if existing_labs or TOKEN else [],
                        metavar="numbers", help=_interval.__doc__)
    parser.add_argument("-s", "--sections", type=int, nargs="+",
                        default=[15, 25], metavar="section",
                        help="your section numbers")


def sheets(labs, sections,
           extensions=["pdf", "png"], force=False):
    """Draw a sign-in sheet showing what group/table students are assigned to.

    Input and output files are organized in the current directory like this:
        .
        ├── canvas.py
        └── lab01
            ├── canvas.csv
            ├── groups015.png
            └── groups025.png
    """
    if verbose:
        print(f"Processing labs [{', '.join(str(lab) for lab in labs)}] and "
              f"sections [{', '.join(str(section) for section in sections)}].")

    for lab in labs:
        # Download CSV if necessary
        canvas_file = os.path.join(f"lab{lab:02d}", "canvas.csv")
        if force or not os.path.exists(canvas_file):
            if verbose:
                print(f"Downloading groups for lab {lab:d} from Canvas.")
            # TODO catch network error to simplify error message
            _canvas_import_csv(lab)
        else:
            if verbose:
                print(f'Using existing file "{canvas_file}".')

        # Parse CSV
        dtype = [("name", "U50"), ("section", int), ("group", int)]
        conv = {0: _format_name,                 # name
                4: lambda name: name[-3:] or 0,  # section
                5: lambda name: name[-1:] or 0}  # group_name
        #cols = list(sorted(conv.keys()))
        names, sect, groups = np.loadtxt(canvas_file,
                                         delimiter=",", quotechar='"',
                                         dtype=dtype, converters=conv,
                                         skiprows=1, usecols=conv.keys(),
                                         unpack=True)

        for section in sections:
            mask = sect == section
            fig = _draw(names[mask], groups[mask],
                        title=f"Groups for Lab {lab:02d} "
                              f"Section {section:03d}")
            for ext in extensions:
                filename = os.path.join(f"lab{lab:02d}",
                                        f"groups{section:03d}.{ext}")
                fig.savefig(filename)
                if verbose:
                    print(f'Output written to "{filename}"')
            plt.close(fig)


def _introduction_parser(parser):
    parser.add_argument("-u", "--update",
                        action=argparse.BooleanOptionalAction,
                        help="update quiz code")
    parser.add_argument("-l", "--lab", type=int,
                        default=max(existing_labs, default=0),
                        metavar="number", help="the lab's number")
    parser.add_argument("-s", "--sections", type=int, nargs="+",
                        default=[15, 25], metavar="section",
                        help="your section numbers")


def introduction(lab, sections, update=False):
    """Create a template for introduction slides

    The template has three slides:
        - title page with lab and first section number
        - sign-in sheets stacked on top of each other
        - a placeholder for the quiz code

    Since the quiz code changes quite frequently we put a placeholder in the
    template. Use the quiz_code command to update it before class.
    """
    intros_path = r"C:\\Users\\umthr\\OneDrive - Umich\\Documents\\Teaching" \
                  r"\\WN26 PHYSICS 251\\Introductions"
    template = intros_path + r"\\Template.pptx"

    if verbose:
        print(f'Using template at "{template}".')
    # Load template presentation
    prs = Presentation(template)

    # Modify title slide
    title_slide = prs.slides[0]
    subtitle = title_slide.placeholders[1]
    subtitle.text = f"Lab {lab:02d} - Section {sections[0]:03d}"

    # Modify group slide
    group_slide = prs.slides[1]
    try:  # looking for an existing picture shape
        pic_shape = next((shape for shape in group_slide.shapes
                          if shape.shape_type == MSO_SHAPE_TYPE.PICTURE), None)
        left   = pic_shape.left
        top    = pic_shape.top
        width  = pic_shape.width
        height = pic_shape.height
        group_slide.shapes._spTree.remove(pic_shape._element)
        if verbose:
            print("Replacing existing picture on group slide.")
    except StopIteration:
        left   = 2487705
        top    = 0
        width  = 6656295
        height = 5143500
        if verbose:
            print("No existing picture found on group slide. "
                  "Falling back to default position and size.")

    # Add sign-in sheets for all sections on top of each other (first on top)
    for section in reversed(sections):
        img_path = f"lab{lab:02d}\\groups{section:03d}.png"
        if not os.path.exists(img_path):
            sheets([lab], [section], extensions=["png"])
        group_slide.shapes.add_picture(img_path, left, top,
                                       width=width, height=height)

    if update:  # quiz code
        quiz_slide = prs.slides[2]
        quiz_code = _get_quiz_code(lab)
        quiz_slide.placeholders[0].text = quiz_code

    # Save the modified presentation
    path = f"{intros_path}\\PHYS251 Lab {lab:02d}.pptx"
    prs.save(path)

    if verbose:
        print(f'Introduction slides for lab {lab:d} saved to "{path}".')


def _get_quiz_code(lab):
    # The quiz code for each lab is defined in a quiz on Canvas
    # see https://developerdocs.instructure.com/services/canvas/resources/quizzes#method.quizzes/quizzes_api.index
    quizzes = _canvas_api(f"courses/{COURSE_ID}/quizzes",
                          parameters={"search_term": f"Quiz {lab:d}:"})
    try:
        quiz = next(quiz for quiz in quizzes
                    if quiz["title"].startswith(f"Quiz {lab:d}:"))
    except StopIteration:
        raise RuntimeError(f"Couldn't find quiz for lab {lab:d} on Canvas.")
    return quiz["access_code"]


def _quiz_code_parser(parser):
    parser.add_argument("-l", "--lab", type=int,
                        default=max(existing_labs, default=0),
                        metavar="number", help="the lab's number")


def quiz_code(lab):
    """Update the quiz code on the introduction slides
    
    This commands pulls the latest quiz code from the Canvas API and updates
    the corresponding slide in the introduction.
    """
    if verbose:
        print(f"Updating quiz code for lab {lab:d}.")

    intros_path = r"C:\\Users\\umthr\\OneDrive - Umich\\Documents\\Teaching" \
                  r"\\WN26 PHYSICS 251\\Introductions"
    intro = f"{intros_path}\\PHYS251 Lab {lab:02d}.pptx"

    if not os.path.exists(intro):
        raise RuntimeError(f'No slides for lab {lab:02d} found at "{intro}".')

    # Load template presentation
    prs = Presentation(intro)

    # Update quiz slide
    quiz_slide = prs.slides[2]
    quiz_code = _get_quiz_code(lab)
    quiz_slide.placeholders[0].text = quiz_code
    if verbose:
        print(f'Quiz code "{quiz_code}" retrieved from Canvas.')

    # Save the modified presentation
    prs.save(intro)
    if verbose:
        print(f'Updated presentation "{intro}".')


def _new_quiz_code_parser(parser):
    parser.add_argument("-l", "--lab", type=int,
                        default=max(existing_labs, default=0),
                        metavar="number", help="the lab's number")


def new_quiz_code(lab):
    """Generate a new quiz code on Canvas
    
    Overwrites the quiz code on Canvas with a random number.
    Intended to be used in class after students finished the quiz.
    """
    if verbose:
        print(f"Generating new quiz code for lab {lab:d}.")

    # The quiz code for each lab is defined in a quiz on Canvas
    # see https://developerdocs.instructure.com/services/canvas/resources/quizzes#method.quizzes-quizzes_api.show
    quizzes = _canvas_api(f"courses/{COURSE_ID}/quizzes",
                          parameters={"search_term": f"Quiz {lab:d}:"})
    try:
        quiz = next(quiz for quiz in quizzes
                    if quiz["title"].startswith(f"Quiz {lab:d}:"))
    except StopIteration:
        raise RuntimeError(f"Couldn't find quiz for lab {lab:d} on Canvas.")

    new_access_code = f"{randrange(10**6):06d}"
    parameters = {"quiz": {"access_code": new_access_code,
                           "notify_of_update": "false"}}
    # see https://developerdocs.instructure.com/services/canvas/resources/quizzes#method.quizzes-quizzes_api.update
    _canvas_api(f"courses/{COURSE_ID}/quizzes/{quiz['id']}",
                method="PUT",
                headers={"Content-Type": "application/json"},
                data=json.dumps(parameters))

    if verbose:
        print(f"Access code for quiz {lab:d} changed to {new_access_code:s}.")


def _get_worksheet(lab, path="Physics 251 GSI Resources/Lab Worksheets"):
    # The worksheet for each lab is defined in a file on Canvas
    # see https://developerdocs.instructure.com/services/canvas/resources/files#method.files.index
    folders = _canvas_api(f"courses/{COURSE_ID}/folders/by_path/{quote(path)}")
    files = _canvas_api(f"folders/{folders[-1]['id']}/files",
                        parameters={"search_term": f"Lab {lab:d} -"})
    if verbose > 1:
        print(f'Found {len(files)} file{'' if len(files) == 1 else 's'} '
              f'in folder "{path}" matching "Lab {lab:d} -".')
        for file in files:
            print(f'  "{file["display_name"]}"',
                  f'unquoted filename: "{unquote_plus(file["filename"])}"',
                  f'id: {file["id"]}',
                  f'url: {file["url"]})', sep="\n    ")
    try:
        file = next(file for file in files
                    if file["display_name"].startswith(f"Lab {lab:d} -"))
    except StopIteration:
        raise RuntimeError("Couldn't find worksheet "
                           f"for lab {lab:d} on Canvas.")
    content = urlopen(file["url"]).read()
    return content


def _worksheet_parser(parser):
    parser.add_argument("-l", "--labs", nargs="+", type=_interval,
                        action=FlatListAction,
                        default=[max(existing_labs, default=0) + 1],
                        metavar="number", help="the lab's number")
    parser.add_argument("-p", "--path", type=str, metavar="path",
                        #default=<parsed from function signature>,
                        help="the path to the worksheet file on Canvas")


def worksheet(labs, path="Physics 251 GSI Resources/Lab Worksheets"):
    """Download the worksheet for the given lab from Canvas

    The worksheets are stored in another Canvas course (at least for PHYS 251
    WN26). Define another Canvas course ID with " GSI" appended to the name and
    update the path to the top folder if necessary.
    """
    if not TOKEN:
        raise RuntimeError("No Canvas API access token defined.")
    
    for lab in labs:
        worksheet = _get_worksheet(lab, path=path)
        with urlopen(worksheet["url"]) as response:
            content = response.read()
        extension = os.path.splitext(worksheet["display_name"])[-1]
        filename = f"lab{lab:02d}\\worksheet.{extension}"
        os.makedirs(os.path.dirname(filename), exist_ok=True)

        with open(filename, "wb") as file:
            file.write(content)
        if verbose:
            print(f'Worksheet for lab {lab:d} written to "{filename}".')


def _letter_grade(grade):
    """FINAL LETTER GRADE CALCULATOR

    PHYSICS 151, 251
    WINTER 2025
    Version 5.2
    Liam Daly, Blake Bottesi, Michelle Thran
    Last updated: 4/22/2025
    """
    if grade >= 99:
        return 'A+'
    if grade >= 94:
        return 'A'
    if grade >= 90:
        return 'A-'
    if grade >= 87:
        return 'B+'
    if grade >= 84:
        return 'B'
    if grade >= 80:
        return 'B-'
    if grade >= 77:
        return 'C+'
    if grade >= 73:
        return 'C'
    if grade >= 70:
        return 'C-'
    if grade >= 67:
        return 'D+'
    if grade >= 63:
        return 'D'
    if grade >= 60:
        return 'D-'
    else:
        return 'E'


def _get_grades(df, grade_key="Final Score"):
    """FINAL LETTER GRADE CALCULATOR

    PHYSICS 151, 251
    WINTER 2025
    Version 5.2
    Liam Daly, Blake Bottesi, Michelle Thran
    Last updated: 4/22/2025
    """
    df['Letter Grade'] = df[grade_key].apply(_letter_grade)
    grades = df[['Student', 'Section', grade_key, 'Letter Grade']]
    return grades


def _uploadable(df):
    """FINAL LETTER GRADE CALCULATOR

    PHYSICS 151, 251
    WINTER 2025
    Version 5.2
    Liam Daly, Blake Bottesi, Michelle Thran
    Last updated: 4/22/2025
    """
    grades = df[['SIS User ID', 'Letter Grade']]
    return grades


def _final_grades_parser(parser):
    parser.add_argument("gradebook",
                        help="path to gradebook exported from Canvas")
    parser.add_argument("-r", "--readable", type=str, metavar="path",
                        #default=<parsed from function signature>,
                        help="file to write human-readable grades")
    parser.add_argument("-u", "--uploadable", type=str, metavar="path",
                        #default=<parsed from function signature>,
                        help="file to write machine-readable grades")
    parser.add_argument("-k", "--grade-key", type=str, metavar="name",
                        #default=<parsed from function signature>,
                        help="the column name for final scores "
                        "in the gradebook CSV")


def final_grades(gradebook, grade_key="Current Score",
                 readable="grades/human-readable.csv",
                 uploadable="grades/wolverine_access.csv"):
    """FINAL LETTER GRADE CALCULATOR

    PHYSICS 151, 251
    WINTER 2025
    Version 5.2
    Liam Daly, Blake Bottesi, Michelle Thran
    Last updated: 4/22/2025
    """
    # Deletes empty rows and test student
    df = pd.read_csv(gradebook, skiprows=[1, 2])
    df = df.drop(df.index[-1])
    if verbose:
        print(f'Successfully read gradebook "{gradebook}" as CSV.')

    # ignores PHYSICS vs BIOPHYS distinction
    df['Section'] = df['Section'].str.slice(-3)

    # Sorts the gradebook first by section, then by student last name
    df = df.sort_values(by=['Section', 'Student'])

    final_grades = _get_grades(df, grade_key=grade_key)

    final_grades.to_csv(readable, index=False)
    if verbose:
        print(f'Human-readable final grades written to "{readable}" as CSV.')

    # Use this csv to cross-reference against canvas to ensure
    # accuracy of letter grades
    uploadable_grades = _uploadable(df)
    uploadable_grades.to_csv(uploadable,
                             index=False, header=None, float_format=int)
    if verbose:
        print(f'Machine-readable file written to "{uploadable}" as CSV.')

    # Use this csv to upload letter grades to wolverine access.
    # Not intended for direct reading


if __name__ == "__main__":
    description, epilog = __doc__.split("\n\n", maxsplit=1)
    parser = argparse.ArgumentParser(description=description, epilog=epilog,
                                     formatter_class=CustomHelpFormatter,
                                     add_help=False)
    parser.add_argument("-h", "--help", action=VerboseHelpAction,
                        help="show this help message and exit")
    parser.add_argument("-v", "--verbose", action="count", default=0,
                        help="print status messages")
    parser.add_argument("-c", "--course", choices=COURSES.keys(),
                        default="PHYS 251 WN26", help="the Canvas course",
                        metavar="name")

    commands = [sheets, introduction, quiz_code, new_quiz_code, worksheet,
                final_grades]
    aliases = {"introduction":  ["intro", "slides"],
               "quiz_code":     ["quiz"],
               "new_quiz_code": ["new_code"],
               "final_grades":  ["grades"]}
    subparsers = parser.add_subparsers(dest="command", metavar="command",
                                       required=True)
    for command in commands:
        description, epilog = command.__doc__.split("\n\n", maxsplit=1)
        subparser = subparsers.add_parser(command.__name__,
                                          aliases=aliases.get(command.__name__,
                                                              []),
                                          formatter_class=CustomHelpFormatter,
                                          description=description,
                                          help=description,
                                          epilog=epilog)
        
        # Construct subparser with command-specific arguments
        if hasattr(command, "parser"):  # e.g., sheets.parser = _sheets_parser
            constructor = command.parser
        else:  # find a constructor by its name
            constructor = globals()[f"_{command.__name__}_parser"]
        constructor(subparser)

        # Read defaults from function signatures
        defaults = {param.name: param.default for param in
                    inspect.signature(command).parameters.values()
                    if param.default is not inspect.Parameter.empty}
        # Set default function to call and its defaults arguments
        subparser.set_defaults(func=command, **defaults)

    args = vars(parser.parse_args())

    # Pop global arguments
    verbose = args.pop("verbose")
    # Find the command to execute
    command_name = args.pop("command")
    command = args.pop("func")
    if verbose:
        if command_name not in [command.__name__ for command in commands]:
            print(f'Command "{command_name}" is an alias '
                  f'for "{command.__name__}".')

    course = args.pop("course")
    if command == worksheet:
        course += " GSI"
    COURSE_ID = COURSES[course]
    if verbose:
        print(f'Using course "{course}" with ID {COURSE_ID:d}.')

    command(**args)